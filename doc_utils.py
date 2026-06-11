"""
Document text extraction utilities.
Supports: PDF, DOCX, PPTX, XLSX, TXT, CSV, MD
All libraries are free/open-source.
"""

import io
import csv


def extract_text_from_file(uploaded_file) -> tuple[str, str]:
    """
    Extract text from a Streamlit UploadedFile object.

    Returns:
        (extracted_text, file_type)

    Raises:
        ValueError if file type is unsupported or extraction fails.
    """
    name = uploaded_file.name.lower()
    file_bytes = uploaded_file.read()

    if name.endswith(".pdf"):
        return _extract_pdf(file_bytes), "PDF"
    elif name.endswith(".docx"):
        return _extract_docx(file_bytes), "Word Document"
    elif name.endswith(".pptx"):
        return _extract_pptx(file_bytes), "PowerPoint"
    elif name.endswith((".xlsx", ".xls")):
        return _extract_excel(file_bytes), "Spreadsheet"
    elif name.endswith((".txt", ".md", ".markdown")):
        return file_bytes.decode("utf-8", errors="replace"), "Text"
    elif name.endswith(".csv"):
        return _extract_csv(file_bytes), "CSV"
    else:
        ext = name.rsplit(".", 1)[-1].upper() if "." in name else "unknown"
        raise ValueError(
            f"Unsupported file type: .{ext}. "
            "Supported: PDF, DOCX, PPTX, XLSX, XLS, TXT, CSV, MD"
        )


def _extract_pdf(file_bytes: bytes) -> str:
    try:
        import pymupdf  # PyMuPDF — faster and more reliable
        doc = pymupdf.open(stream=file_bytes, filetype="pdf")
        pages = []
        for page in doc:
            pages.append(page.get_text())
        return "\n\n".join(pages)
    except ImportError:
        pass

    # Fallback to pypdf
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        return "\n\n".join(
            page.extract_text() or "" for page in reader.pages
        )
    except ImportError:
        raise ImportError(
            "PDF extraction requires pymupdf or pypdf. "
            "Install with: pip install pymupdf"
        )


def _extract_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "DOCX extraction requires python-docx. "
            "Install with: pip install python-docx"
        )
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract tables
    table_texts = []
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_texts.append(" | ".join(cells))
    all_text = paragraphs + table_texts
    return "\n\n".join(all_text)


def _extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "PPTX extraction requires python-pptx. "
            "Install with: pip install python-pptx"
        )
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if len(slide_lines) > 1:
            slides_text.append("\n".join(slide_lines))
    return "\n\n".join(slides_text)


def _extract_excel(file_bytes: bytes) -> str:
    try:
        import openpyxl
    except ImportError:
        raise ImportError(
            "Excel extraction requires openpyxl. "
            "Install with: pip install openpyxl"
        )
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    return "\n\n".join(sheets_text)


def _extract_csv(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        if any(cell.strip() for cell in row):
            rows.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(rows)


def get_file_icon(file_type: str) -> str:
    icons = {
        "PDF": "📄",
        "Word Document": "📝",
        "PowerPoint": "📊",
        "Spreadsheet": "📈",
        "Text": "📃",
        "CSV": "📋",
    }
    return icons.get(file_type, "📁")
